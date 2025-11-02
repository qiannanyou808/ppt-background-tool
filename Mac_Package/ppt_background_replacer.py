#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT背景替换工具
支持批量替换PPT文件的背景图片
"""

import os
import threading
from pathlib import Path
from tkinter import filedialog, messagebox
import customtkinter as ctk
from pptx import Presentation

# 设置主题
ctk.set_appearance_mode("light")
ctk.set_default_color_theme("blue")


class PPTBackgroundReplacer(ctk.CTk):
    def __init__(self):
        super().__init__()

        # 窗口配置
        self.title("PPT背景替换工具")
        self.geometry("800x700")
        self.resizable(False, False)

        # 居中显示
        self.center_window()

        # 变量
        self.ppt_files = []
        self.background_image = None
        self.output_dir = None  # 输出目录（None表示保存到原文件目录）
        self.is_processing = False

        # 创建UI
        self.create_widgets()

    def center_window(self):
        """窗口居中显示"""
        self.update_idletasks()
        width = self.winfo_width()
        height = self.winfo_height()
        x = (self.winfo_screenwidth() // 2) - (width // 2)
        y = (self.winfo_screenheight() // 2) - (height // 2)
        self.geometry(f'{width}x{height}+{x}+{y}')

    def create_widgets(self):
        """创建界面组件"""
        # 标题区域
        title_frame = ctk.CTkFrame(self, fg_color="transparent")
        title_frame.pack(pady=20, padx=20, fill="x")

        title_label = ctk.CTkLabel(
            title_frame,
            text="📊 PPT背景替换工具",
            font=ctk.CTkFont(size=28, weight="bold")
        )
        title_label.pack()

        subtitle_label = ctk.CTkLabel(
            title_frame,
            text="批量替换PPT背景，保留所有内容",
            font=ctk.CTkFont(size=14),
            text_color="gray"
        )
        subtitle_label.pack(pady=(5, 0))

        # 主要内容区域
        main_frame = ctk.CTkFrame(self)
        main_frame.pack(pady=10, padx=40, fill="both", expand=True)

        # 步骤1: 选择PPT文件
        step1_frame = ctk.CTkFrame(main_frame)
        step1_frame.pack(pady=8, padx=20, fill="x")

        step1_label = ctk.CTkLabel(
            step1_frame,
            text="步骤 1: 选择PPT文件",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step1_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame1 = ctk.CTkFrame(step1_frame, fg_color="transparent")
        btn_frame1.pack(padx=10, pady=(0, 10), fill="x")

        self.btn_select_ppt = ctk.CTkButton(
            btn_frame1,
            text="选择PPT文件",
            command=self.select_ppt_files,
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_ppt.pack(side="left")

        self.label_ppt_count = ctk.CTkLabel(
            btn_frame1,
            text="未选择文件",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_ppt_count.pack(side="left", padx=15)

        # 步骤2: 选择背景图片
        step2_frame = ctk.CTkFrame(main_frame)
        step2_frame.pack(pady=8, padx=20, fill="x")

        step2_label = ctk.CTkLabel(
            step2_frame,
            text="步骤 2: 选择背景图片",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step2_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame2 = ctk.CTkFrame(step2_frame, fg_color="transparent")
        btn_frame2.pack(padx=10, pady=(0, 10), fill="x")

        self.btn_select_bg = ctk.CTkButton(
            btn_frame2,
            text="选择背景图片",
            command=self.select_background_image,
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_bg.pack(side="left")

        self.label_bg_path = ctk.CTkLabel(
            btn_frame2,
            text="未选择图片",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_bg_path.pack(side="left", padx=15)

        # 步骤3: 选择输出目录（可选）
        step3_frame = ctk.CTkFrame(main_frame)
        step3_frame.pack(pady=8, padx=20, fill="x")

        step3_label = ctk.CTkLabel(
            step3_frame,
            text="步骤 3: 选择输出目录（可选）",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step3_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame3 = ctk.CTkFrame(step3_frame, fg_color="transparent")
        btn_frame3.pack(padx=10, pady=(0, 10), fill="x")

        self.btn_select_output = ctk.CTkButton(
            btn_frame3,
            text="选择输出目录",
            command=self.select_output_directory,
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_output.pack(side="left")

        self.label_output_dir = ctk.CTkLabel(
            btn_frame3,
            text="默认：保存到原文件目录",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_output_dir.pack(side="left", padx=15)

        # 步骤4: 开始替换
        step4_frame = ctk.CTkFrame(main_frame)
        step4_frame.pack(pady=8, padx=20, fill="x")

        step4_label = ctk.CTkLabel(
            step4_frame,
            text="步骤 4: 开始处理",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step4_label.pack(anchor="w", padx=10, pady=(10, 5))

        self.btn_process = ctk.CTkButton(
            step4_frame,
            text="🚀 开始替换背景",
            command=self.start_processing,
            width=200,
            height=45,
            font=ctk.CTkFont(size=16, weight="bold"),
            fg_color="#2B8A3E",
            hover_color="#1D5C29"
        )
        self.btn_process.pack(padx=10, pady=(0, 10))

        # 进度显示区域
        progress_frame = ctk.CTkFrame(main_frame)
        progress_frame.pack(pady=10, padx=20, fill="both", expand=True)

        progress_title = ctk.CTkLabel(
            progress_frame,
            text="处理进度",
            font=ctk.CTkFont(size=14, weight="bold")
        )
        progress_title.pack(anchor="w", padx=15, pady=(15, 10))

        # 当前处理文件名显示
        self.current_file_label = ctk.CTkLabel(
            progress_frame,
            text="等待开始...",
            font=ctk.CTkFont(size=14, weight="bold"),
            text_color="#1F6AA5"
        )
        self.current_file_label.pack(pady=(5, 10), padx=15)

        # 进度百分比和状态
        self.progress_label = ctk.CTkLabel(
            progress_frame,
            text="准备就绪",
            font=ctk.CTkFont(size=12),
            text_color="gray"
        )
        self.progress_label.pack(pady=(0, 8), padx=15)

        # 进度条
        self.progress_bar = ctk.CTkProgressBar(
            progress_frame,
            height=25
        )
        self.progress_bar.pack(fill="x", padx=15, pady=(0, 8))
        self.progress_bar.set(0)

        # 页面处理进度显示
        self.page_progress_label = ctk.CTkLabel(
            progress_frame,
            text="",
            font=ctk.CTkFont(size=11),
            text_color="gray"
        )
        self.page_progress_label.pack(pady=(0, 10), padx=15)

        # 简洁日志显示（只显示最新的3-5条）
        log_label = ctk.CTkLabel(
            progress_frame,
            text="处理日志",
            font=ctk.CTkFont(size=12, weight="bold")
        )
        log_label.pack(anchor="w", padx=15, pady=(10, 5))

        self.status_text = ctk.CTkTextbox(
            progress_frame,
            height=120,
            font=ctk.CTkFont(size=10),
            wrap="none"
        )
        self.status_text.pack(
            padx=15, pady=(0, 15), fill="both", expand=True
        )
        self.status_text.insert("1.0", "等待操作...\n")
        self.status_text.configure(state="disabled")

        # 底部信息
        footer_label = ctk.CTkLabel(
            self,
            text="提示：处理后的文件将保存为 原文件名_新背景.pptx",
            font=ctk.CTkFont(size=11),
            text_color="gray"
        )
        footer_label.pack(pady=(0, 10))

    def select_ppt_files(self):
        """选择PPT文件"""
        files = filedialog.askopenfilenames(
            title="选择PPT文件",
            filetypes=[
                ("PowerPoint文件", "*.pptx"),
                ("所有文件", "*.*")
            ]
        )

        if files:
            self.ppt_files = list(files)
            count = len(self.ppt_files)
            self.label_ppt_count.configure(
                text=f"已选择 {count} 个文件",
                text_color="#2B8A3E"
            )
            self.log_status(f"✓ 已选择 {count} 个PPT文件")

    def select_background_image(self):
        """选择背景图片"""
        file = filedialog.askopenfilename(
            title="选择背景图片",
            filetypes=[
                ("图片文件", "*.jpg *.jpeg *.png"),
                ("所有文件", "*.*")
            ]
        )

        if file:
            self.background_image = file
            filename = os.path.basename(file)
            # 截断过长的文件名
            if len(filename) > 30:
                filename = filename[:27] + "..."
            self.label_bg_path.configure(
                text=filename,
                text_color="#2B8A3E"
            )
            bg_name = os.path.basename(file)
            self.log_status(f"✓ 已选择背景图片: {bg_name}")

    def select_output_directory(self):
        """选择输出目录"""
        directory = filedialog.askdirectory(
            title="选择输出目录"
        )

        if directory:
            self.output_dir = directory
            # 显示目录路径（如果太长则截断）
            display_path = directory
            if len(display_path) > 40:
                display_path = "..." + display_path[-37:]
            self.label_output_dir.configure(
                text=f"输出到: {display_path}",
                text_color="#2B8A3E"
            )
            self.log_status(f"✓ 已选择输出目录: {directory}")

    def log_status(self, message):
        """记录状态信息"""
        # 在GUI中显示
        self.status_text.configure(state="normal")
        self.status_text.insert("end", f"{message}\n")
        self.status_text.see("end")
        self.status_text.configure(state="disabled")

        # 同时在终端输出（方便查看详细日志）
        try:
            print(message)
        except UnicodeEncodeError:
            # Windows终端可能不支持某些Unicode字符，降级处理
            safe_message = (
                message.encode('gbk', errors='replace').decode('gbk')
            )
            print(safe_message)

    def start_processing(self):
        """开始处理"""
        if self.is_processing:
            messagebox.showwarning("提示", "正在处理中，请稍候...")
            return

        if not self.ppt_files:
            messagebox.showwarning("提示", "请先选择PPT文件！")
            return

        if not self.background_image:
            messagebox.showwarning("提示", "请先选择背景图片！")
            return

        # 在新线程中处理，避免界面卡顿
        thread = threading.Thread(target=self.process_ppts, daemon=True)
        thread.start()

    def process_ppts(self):
        """处理PPT文件"""
        self.is_processing = True
        self.btn_process.configure(state="disabled")

        total_files = len(self.ppt_files)
        success_count = 0
        fail_count = 0
        last_output_file = None  # 记录最后一个成功的输出文件

        self.log_status("\n" + "="*50)
        self.log_status("开始处理...")
        self.log_status("="*50)

        for idx, ppt_file in enumerate(self.ppt_files, 1):
            try:
                # 获取文件名
                filename = os.path.basename(ppt_file)

                # 更新当前处理文件显示
                self.current_file_label.configure(
                    text=f"📄 正在处理: {filename}",
                    text_color="#1F6AA5"
                )

                # 更新进度
                progress = (idx - 1) / total_files
                self.progress_bar.set(progress)
                self.progress_label.configure(
                    text=f"文件进度: {idx-1}/{total_files} "
                         f"({progress*100:.0f}%)"
                )

                # 清空页面进度显示
                self.page_progress_label.configure(text="")

                self.log_status(
                    f"\n[{idx}/{total_files}] 正在处理: {filename}"
                )

                # 处理单个PPT
                output_file = self.replace_background(ppt_file)

                output_name = os.path.basename(output_file)
                self.log_status(f"✓ 处理成功: {output_name}")
                success_count += 1
                last_output_file = output_file  # 记录成功的文件

            except Exception as e:
                import traceback
                error_detail = traceback.format_exc()
                self.log_status(f"✗ 处理失败: {str(e)}")
                self.log_status(f"详细错误:\n{error_detail}")
                fail_count += 1

        # 完成
        self.progress_bar.set(1.0)
        self.current_file_label.configure(
            text="✅ 全部处理完成！",
            text_color="#2B8A3E"
        )
        self.progress_label.configure(
            text=f"✓ 完成: {total_files}/{total_files} (100%)",
            text_color="#2B8A3E"
        )
        self.page_progress_label.configure(
            text=f"成功: {success_count} 个 | 失败: {fail_count} 个",
            text_color="#2B8A3E" if fail_count == 0 else "#C92A2A"
        )
        self.log_status("\n" + "="*50)
        status_msg = f"处理完成！成功: {success_count}, 失败: {fail_count}"
        self.log_status(status_msg)
        self.log_status("="*50 + "\n")

        self.is_processing = False
        self.btn_process.configure(state="normal")

        # 显示完成提示
        result_msg = (
            f"处理完成！\n\n"
            f"成功: {success_count} 个文件\n"
            f"失败: {fail_count} 个文件"
        )
        messagebox.showinfo("处理完成", result_msg)

        # 如果有成功的文件，打开文件所在位置
        if last_output_file and success_count > 0:
            try:
                import subprocess
                # 使用Windows资源管理器打开并选中文件
                subprocess.Popen(
                    f'explorer /select,"{last_output_file}"'
                )
                self.log_status("✓ 已打开文件所在位置")
            except Exception as e:
                self.log_status(f"无法打开文件位置: {str(e)}")

    def replace_background(self, ppt_file):
        """替换单个PPT的背景 - 使用简单可靠的方法"""
        from pptx.util import Inches

        self.log_status("  → 打开PPT文件...")
        # 打开PPT
        prs = Presentation(ppt_file)
        self.log_status(f"  → PPT共有 {len(prs.slides)} 页")

        # 获取幻灯片尺寸
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 为每一页设置背景
        for idx, slide in enumerate(prs.slides, 1):
            # 更新页面进度显示
            page_progress = idx / len(prs.slides)
            self.page_progress_label.configure(
                text=f"页面进度: {idx}/{len(prs.slides)} "
                     f"({page_progress*100:.0f}%)"
            )

            self.log_status(f"  → 处理第 {idx}/{len(prs.slides)} 页...")

            # 步骤1：移除原有背景
            try:
                # 1.1 删除背景元素
                slide_elem = slide.element
                namespaces = {
                    'p': 'http://schemas.openxmlformats.org/'
                         'presentationml/2006/main',
                    'a': 'http://schemas.openxmlformats.org/'
                         'drawingml/2006/main'
                }

                bg_removed = False
                for bg in slide_elem.findall('.//p:bg', namespaces):
                    bg.getparent().remove(bg)
                    bg_removed = True

                if bg_removed:
                    self.log_status("    ✓ 已清除背景填充")

                # 1.2 隐藏母版背景
                cSld = slide_elem.find('.//p:cSld', namespaces)
                if cSld is not None:
                    cSld.set('showMasterSp', '0')
                    self.log_status("    ✓ 已隐藏母版背景")

                # 1.3 删除看起来像背景的图片和形状
                # 识别标准（优化后，包括渐变蒙版）：
                # - 完全铺满的图片/形状
                # - 横向撑满的图片/形状（降低面积要求到20%）
                # - 纵向撑满的图片/形状（降低面积要求到20%）
                # - 宽度或高度占比>85%的接近撑满图片/形状
                # - 面积占比超过60%的大图片/形状
                shapes_to_remove = []
                tolerance = Inches(0.15)  # 允许的误差范围

                # 计算幻灯片面积
                slide_area = slide_width * slide_height

                for shape in slide.shapes:
                    try:
                        # 跳过占位符（标题框、内容框等）
                        if (hasattr(shape, 'is_placeholder') and
                                shape.is_placeholder):
                            continue

                        # 检查形状类型：
                        # 1 = 自动形状（包括渐变蒙版）
                        # 13 = 图片
                        # 6 = 组合形状
                        if shape.shape_type in [1, 13, 6]:
                            # 计算图片面积和尺寸占比
                            shape_area = shape.width * shape.height
                            area_ratio = shape_area / slide_area
                            width_ratio = shape.width / slide_width
                            height_ratio = shape.height / slide_height

                            # 检查是否横向撑满（严格匹配）
                            is_width_full = (
                                abs(shape.width - slide_width) < tolerance
                            )

                            # 检查是否纵向撑满（严格匹配）
                            is_height_full = (
                                abs(shape.height - slide_height) < tolerance
                            )

                            # 检查是否接近撑满（多个级别）
                            is_width_near_full = width_ratio > 0.85  # 非常接近
                            is_height_near_full = height_ratio > 0.85
                            is_width_large = width_ratio > 0.70  # 较大
                            is_height_large = height_ratio > 0.70

                            # 检查是否完全铺满
                            is_full_size = is_width_full and is_height_full

                            # 检查面积占比（多级阈值）
                            is_large_area = area_ratio > 0.6
                            is_medium_area = area_ratio > 0.4
                            is_small_area = area_ratio > 0.15

                            # 判断是否为背景图片（优化后的规则）
                            is_background = (
                                is_full_size or  # 完全铺满
                                # 横向撑满且占比>20%（降低阈值）
                                (is_width_full and area_ratio > 0.2) or
                                # 纵向撑满且占比>20%（降低阈值）
                                (is_height_full and area_ratio > 0.2) or
                                # 宽度接近撑满(>85%)且占比>30%
                                (is_width_near_full and area_ratio > 0.3) or
                                # 高度接近撑满(>85%)且占比>30%
                                (is_height_near_full and area_ratio > 0.3) or
                                # 宽度较大(>70%)且占比>40%（新增）
                                (is_width_large and is_medium_area) or
                                # 高度较大(>70%)且占比>40%（新增）
                                (is_height_large and is_medium_area) or
                                # 宽度很大(>85%)且占比>15%（新增，捕获窄装饰条）
                                (is_width_near_full and is_small_area) or
                                # 高度很大(>85%)且占比>15%（新增，捕获窄装饰条）
                                (is_height_near_full and is_small_area) or
                                # 面积占比>60%
                                is_large_area
                            )

                            if is_background:
                                shapes_to_remove.append(shape)

                                # 判断形状类型
                                shape_type_name = "图片"
                                if shape.shape_type == 1:
                                    shape_type_name = "形状"
                                elif shape.shape_type == 6:
                                    shape_type_name = "组合"

                                reason = ""
                                if is_full_size:
                                    reason = "完全铺满"
                                elif is_width_full:
                                    reason = f"横向撑满({area_ratio*100:.0f}%)"
                                elif is_height_full:
                                    reason = f"纵向撑满({area_ratio*100:.0f}%)"
                                elif is_width_near_full or is_width_large:
                                    reason = (
                                        f"宽度{width_ratio*100:.0f}%"
                                        f"/面积{area_ratio*100:.0f}%"
                                    )
                                elif is_height_near_full or is_height_large:
                                    reason = (
                                        f"高度{height_ratio*100:.0f}%"
                                        f"/面积{area_ratio*100:.0f}%"
                                    )
                                else:
                                    reason = f"面积{area_ratio*100:.0f}%"

                                self.log_status(
                                    f"    → 发现背景{shape_type_name}: {reason} "
                                    f"({shape.width/914400:.1f}x"
                                    f"{shape.height/914400:.1f}英寸)"
                                )
                    except Exception:
                        continue

                # 删除识别出的背景图片和形状
                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

                if shapes_to_remove:
                    self.log_status(
                        f"    ✓ 已删除 {len(shapes_to_remove)} 个背景元素"
                    )

            except Exception as e:
                self.log_status(f"    - 处理背景: {str(e)[:50]}")

            # 步骤2：在最底层添加新的背景图片
            left = Inches(0)
            top = Inches(0)

            # 添加图片
            pic = slide.shapes.add_picture(
                self.background_image,
                left, top,
                width=slide_width,
                height=slide_height
            )

            # 将图片移到最底层
            # 获取shape tree
            shapes = slide.shapes._spTree
            # 移除刚添加的图片
            shapes.remove(pic._element)
            # 插入到索引2的位置（索引0和1通常是固定元素）
            shapes.insert(2, pic._element)

        self.log_status("  → 生成输出文件...")
        # 生成输出文件名
        file_path = Path(ppt_file)
        new_name = f"{file_path.stem}_新背景{file_path.suffix}"

        # 根据用户选择的输出目录决定保存位置
        if self.output_dir:
            # 用户指定了输出目录
            output_file = Path(self.output_dir) / new_name
        else:
            # 默认：保存到原文件目录
            output_file = file_path.parent / new_name

        # 如果文件已存在，添加数字后缀
        counter = 1
        while output_file.exists():
            new_name = f"{file_path.stem}_新背景_{counter}{file_path.suffix}"
            if self.output_dir:
                output_file = Path(self.output_dir) / new_name
            else:
                output_file = file_path.parent / new_name
            counter += 1

        self.log_status("  → 保存文件...")
        # 保存
        prs.save(str(output_file))
        self.log_status(f"  ✓ 已保存: {output_file.name}")

        return str(output_file)


def main():
    """主函数"""
    app = PPTBackgroundReplacer()
    app.mainloop()


if __name__ == "__main__":
    main()
