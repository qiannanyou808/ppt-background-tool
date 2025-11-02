#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT背景替换工具 - 学习版（带详细注释）

这个文件是专门为学习者准备的，包含详细的中文注释
建议配合《新手学习指南.txt》一起学习
"""

# ============================================================
# 第一部分：导入需要的模块（相当于导入工具包）
# ============================================================

import os  # 操作系统相关功能，如文件路径处理
import threading  # 多线程，用于避免界面卡死
from pathlib import Path  # 更现代的路径处理方式
from tkinter import filedialog, messagebox  # 文件对话框和消息框
import customtkinter as ctk  # 现代化的GUI界面库
from pptx import Presentation  # 处理PowerPoint文件的库

# ============================================================
# 设置界面主题
# ============================================================
ctk.set_appearance_mode("light")  # 设置为亮色主题
ctk.set_default_color_theme("blue")  # 设置主题色为蓝色


# ============================================================
# 第二部分：主类定义
# ============================================================

class PPTBackgroundReplacer(ctk.CTk):
    """
    PPT背景替换工具的主类
    
    继承自 ctk.CTk（CustomTkinter的窗口类）
    这个类包含了程序的所有功能
    """
    
    def __init__(self):
        """
        初始化方法（构造函数）
        
        当创建这个类的实例时，这个方法会自动执行
        用于设置窗口的基本属性和初始化变量
        """
        # 调用父类的初始化方法（必须的）
        super().__init__()

        # 窗口配置
        self.title("PPT背景替换工具")  # 设置窗口标题
        self.geometry("800x700")  # 设置窗口大小：宽800像素，高700像素（增加高度以确保进度条可见）
        self.resizable(False, False)  # 禁止调整窗口大小

        # 居中显示窗口
        self.center_window()

        # 初始化变量（存储程序运行时需要的数据）
        self.ppt_files = []  # 用户选择的PPT文件列表（空列表）
        self.background_image = None  # 用户选择的背景图片路径（初始为None）
        self.output_dir = None  # 输出目录（None表示保存到原文件目录）
        self.is_processing = False  # 是否正在处理（防止重复点击）

        # 创建界面组件
        self.create_widgets()

    def center_window(self):
        """
        让窗口在屏幕中央显示
        
        计算屏幕中心位置，然后移动窗口
        """
        self.update_idletasks()  # 更新窗口信息
        width = self.winfo_width()  # 获取窗口宽度
        height = self.winfo_height()  # 获取窗口高度
        x = (self.winfo_screenwidth() // 2) - (width // 2)  # 计算X坐标
        y = (self.winfo_screenheight() // 2) - (height // 2)  # 计算Y坐标
        self.geometry(f'{width}x{height}+{x}+{y}')  # 设置新位置

    def create_widgets(self):
        """
        创建界面组件
        
        这是最长的方法，创建了所有的按钮、标签、文本框等
        使用CTkLabel、CTkButton、CTkTextbox等组件
        """
        # ========== 标题区域 ==========
        title_frame = ctk.CTkFrame(self, fg_color="transparent")
        title_frame.pack(pady=20, padx=20, fill="x")

        # 主标题（带emoji图标）
        title_label = ctk.CTkLabel(
            title_frame,
            text="📊 PPT背景替换工具",
            font=ctk.CTkFont(size=28, weight="bold")  # 大号粗体字
        )
        title_label.pack()

        # 副标题（说明文字）
        subtitle_label = ctk.CTkLabel(
            title_frame,
            text="批量替换PPT背景，保留所有内容",
            font=ctk.CTkFont(size=14),
            text_color="gray"  # 灰色文字
        )
        subtitle_label.pack(pady=(5, 0))

        # ========== 主要内容区域 ==========
        main_frame = ctk.CTkFrame(self)
        main_frame.pack(pady=10, padx=40, fill="both", expand=True)

        # ========== 步骤1: 选择PPT文件 ==========
        step1_frame = ctk.CTkFrame(main_frame)
        step1_frame.pack(pady=8, padx=20, fill="x")

        step1_label = ctk.CTkLabel(
            step1_frame,
            text="步骤 1: 选择PPT文件",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step1_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame1 = ctk.CTkFrame(step1_frame, fg_color="transparent")
        btn_frame1.pack(padx=10, pady=(0, 10), fill="x")

        # "选择PPT文件"按钮
        self.btn_select_ppt = ctk.CTkButton(
            btn_frame1,
            text="选择PPT文件",
            command=self.select_ppt_files,  # 点击时调用select_ppt_files方法
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_ppt.pack(side="left")

        # 显示已选择文件数量的标签
        self.label_ppt_count = ctk.CTkLabel(
            btn_frame1,
            text="未选择文件",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_ppt_count.pack(side="left", padx=15)

        # ========== 步骤2: 选择背景图片 ==========
        # （类似步骤1，省略注释）
        step2_frame = ctk.CTkFrame(main_frame)
        step2_frame.pack(pady=8, padx=20, fill="x")

        step2_label = ctk.CTkLabel(
            step2_frame,
            text="步骤 2: 选择背景图片",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step2_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame2 = ctk.CTkFrame(step2_frame, fg_color="transparent")
        btn_frame2.pack(padx=10, pady=(0, 10), fill="x")

        self.btn_select_bg = ctk.CTkButton(
            btn_frame2,
            text="选择背景图片",
            command=self.select_background_image,
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_bg.pack(side="left")

        self.label_bg_path = ctk.CTkLabel(
            btn_frame2,
            text="未选择图片",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_bg_path.pack(side="left", padx=15)

        # ========== 步骤3: 选择输出目录（可选）==========
        step3_frame = ctk.CTkFrame(main_frame)
        step3_frame.pack(pady=8, padx=20, fill="x")

        step3_label = ctk.CTkLabel(
            step3_frame,
            text="步骤 3: 选择输出目录（可选）",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step3_label.pack(anchor="w", padx=10, pady=(10, 5))

        btn_frame3 = ctk.CTkFrame(step3_frame, fg_color="transparent")
        btn_frame3.pack(padx=10, pady=(0, 10), fill="x")

        # "选择输出目录"按钮
        self.btn_select_output = ctk.CTkButton(
            btn_frame3,
            text="选择输出目录",
            command=self.select_output_directory,
            width=150,
            height=35,
            font=ctk.CTkFont(size=14)
        )
        self.btn_select_output.pack(side="left")

        # 显示输出目录路径的标签
        self.label_output_dir = ctk.CTkLabel(
            btn_frame3,
            text="默认：保存到原文件目录",
            font=ctk.CTkFont(size=13),
            text_color="gray"
        )
        self.label_output_dir.pack(side="left", padx=15)

        # ========== 步骤4: 开始替换 ==========
        step4_frame = ctk.CTkFrame(main_frame)
        step4_frame.pack(pady=8, padx=20, fill="x")

        step4_label = ctk.CTkLabel(
            step4_frame,
            text="步骤 4: 开始处理",
            font=ctk.CTkFont(size=16, weight="bold")
        )
        step4_label.pack(anchor="w", padx=10, pady=(10, 5))

        # "开始替换背景"按钮（绿色）
        self.btn_process = ctk.CTkButton(
            step4_frame,
            text="🚀 开始替换背景",
            command=self.start_processing,
            width=200,
            height=45,
            font=ctk.CTkFont(size=16, weight="bold"),
            fg_color="#2B8A3E",  # 绿色背景
            hover_color="#1D5C29"  # 鼠标悬停时的颜色
        )
        self.btn_process.pack(padx=10, pady=(0, 10))

        # ========== 进度显示区域 ==========
        # 这个区域包含进度条和简要日志
        progress_frame = ctk.CTkFrame(main_frame)
        progress_frame.pack(pady=10, padx=20, fill="both", expand=True)

        # 区域标题
        progress_title = ctk.CTkLabel(
            progress_frame,
            text="处理进度",
            font=ctk.CTkFont(size=14, weight="bold")
        )
        progress_title.pack(anchor="w", padx=15, pady=(15, 10))

        # 当前处理文件名显示（大字体，醒目）
        self.current_file_label = ctk.CTkLabel(
            progress_frame,
            text="等待开始...",
            font=ctk.CTkFont(size=14, weight="bold"),
            text_color="#1F6AA5"  # 蓝色
        )
        self.current_file_label.pack(pady=(5, 10), padx=15)

        # 进度百分比和状态显示
        self.progress_label = ctk.CTkLabel(
            progress_frame,
            text="准备就绪",
            font=ctk.CTkFont(size=12),
            text_color="gray"
        )
        self.progress_label.pack(pady=(0, 8), padx=15)

        # 进度条（高度增加到25像素）
        self.progress_bar = ctk.CTkProgressBar(
            progress_frame,
            height=25
        )
        self.progress_bar.pack(fill="x", padx=15, pady=(0, 8))
        self.progress_bar.set(0)  # 初始值为0（0%）

        # 页面处理进度显示（显示正在处理第几页）
        self.page_progress_label = ctk.CTkLabel(
            progress_frame,
            text="",
            font=ctk.CTkFont(size=11),
            text_color="gray"
        )
        self.page_progress_label.pack(pady=(0, 10), padx=15)

        # 简洁日志显示（高度120像素，只显示最新几条）
        log_label = ctk.CTkLabel(
            progress_frame,
            text="处理日志",
            font=ctk.CTkFont(size=12, weight="bold")
        )
        log_label.pack(anchor="w", padx=15, pady=(10, 5))

        # 文本框：显示简要处理日志
        self.status_text = ctk.CTkTextbox(
            progress_frame,
            height=120,  # 缩小日志区域
            font=ctk.CTkFont(size=10),
            wrap="none"
        )
        self.status_text.pack(
            padx=15, pady=(0, 15), fill="both", expand=True
        )
        # 插入初始文本
        self.status_text.insert("1.0", "等待操作...\n")
        # 设置为只读（用户不能编辑）
        self.status_text.configure(state="disabled")

        # ========== 底部提示 ==========
        footer_label = ctk.CTkLabel(
            self,
            text="提示：处理后的文件将保存为 原文件名_新背景.pptx",
            font=ctk.CTkFont(size=11),
            text_color="gray"
        )
        footer_label.pack(pady=(0, 10))

    # ============================================================
    # 第三部分：文件选择功能
    # ============================================================

    def select_ppt_files(self):
        """
        选择PPT文件
        
        弹出文件选择对话框，让用户选择一个或多个PPT文件
        """
        # 调用文件选择对话框
        files = filedialog.askopenfilenames(
            title="选择PPT文件",  # 对话框标题
            filetypes=[  # 文件类型过滤
                ("PowerPoint文件", "*.pptx"),  # 只显示.pptx文件
                ("所有文件", "*.*")  # 也可以显示所有文件
            ]
        )

        # 如果用户选择了文件（不是取消）
        if files:
            # 将选择的文件保存到self.ppt_files
            self.ppt_files = list(files)
            
            # 统计文件数量
            count = len(self.ppt_files)
            
            # 更新界面上的提示文字
            self.label_ppt_count.configure(
                text=f"已选择 {count} 个文件",
                text_color="#2B8A3E"  # 绿色表示成功
            )
            
            # 在状态框显示日志
            self.log_status(f"✓ 已选择 {count} 个PPT文件")

    def select_background_image(self):
        """
        选择背景图片
        
        弹出文件选择对话框，让用户选择一张图片
        """
        # 调用文件选择对话框（单选）
        file = filedialog.askopenfilename(
            title="选择背景图片",
            filetypes=[
                ("图片文件", "*.jpg *.jpeg *.png"),  # 支持的图片格式
                ("所有文件", "*.*")
            ]
        )

        # 如果用户选择了文件
        if file:
            # 保存图片路径
            self.background_image = file
            
            # 获取文件名（不含路径）
            filename = os.path.basename(file)
            
            # 如果文件名太长，截断显示
            if len(filename) > 30:
                filename = filename[:27] + "..."
            
            # 更新界面显示
            self.label_bg_path.configure(
                text=filename,
                text_color="#2B8A3E"
            )
            
            # 记录日志
            bg_name = os.path.basename(file)
            self.log_status(f"✓ 已选择背景图片: {bg_name}")

    def select_output_directory(self):
        """
        选择输出目录（可选）
        
        弹出目录选择对话框，让用户选择一个文件夹
        如果不选择，则默认保存到原文件所在目录
        """
        # 调用目录选择对话框
        directory = filedialog.askdirectory(
            title="选择输出目录"
        )

        # 如果用户选择了目录
        if directory:
            # 保存输出目录路径
            self.output_dir = directory

            # 显示目录路径（如果太长则截断）
            display_path = directory
            if len(display_path) > 40:
                # 只显示路径的后37个字符，前面加...
                display_path = "..." + display_path[-37:]

            # 更新界面显示
            self.label_output_dir.configure(
                text=f"输出到: {display_path}",
                text_color="#2B8A3E"  # 绿色表示成功
            )

            # 记录日志
            self.log_status(f"✓ 已选择输出目录: {directory}")

    # ============================================================
    # 第四部分：日志显示功能
    # ============================================================

    def log_status(self, message):
        """
        在状态框显示日志消息
        
        参数:
            message: 要显示的消息文本
        """
        # 设置文本框为可编辑
        self.status_text.configure(state="normal")

        # 在末尾插入新消息
        self.status_text.insert("end", f"{message}\n")

        # 滚动到最新的消息
        self.status_text.see("end")

        # 设置回只读状态
        self.status_text.configure(state="disabled")

        # 同时在终端输出（方便查看详细日志）
        try:
            print(message)
        except UnicodeEncodeError:
            # Windows终端可能不支持某些Unicode字符，降级处理
            safe_message = (
                message.encode('gbk', errors='replace').decode('gbk')
            )
            print(safe_message)

    # ============================================================
    # 第五部分：处理流程控制
    # ============================================================

    def start_processing(self):
        """
        开始处理（点击"开始替换背景"按钮时调用）
        
        进行检查后，创建新线程开始处理
        """
        # 检查：是否正在处理
        if self.is_processing:
            messagebox.showwarning("提示", "正在处理中，请稍候...")
            return

        # 检查：是否选择了PPT文件
        if not self.ppt_files:
            messagebox.showwarning("提示", "请先选择PPT文件！")
            return

        # 检查：是否选择了背景图片
        if not self.background_image:
            messagebox.showwarning("提示", "请先选择背景图片！")
            return

        # 在新线程中处理，避免界面卡死
        # daemon=True 表示主程序退出时，这个线程也会自动结束
        thread = threading.Thread(
            target=self.process_ppts,  # 目标函数
            daemon=True
        )
        thread.start()  # 启动线程

    # ============================================================
    # 第六部分：批量处理逻辑
    # ============================================================

    def process_ppts(self):
        """
        批量处理PPT文件
        
        遍历用户选择的所有PPT文件，逐个处理
        """
        # 标记为正在处理
        self.is_processing = True
        # 禁用"开始替换背景"按钮（防止重复点击）
        self.btn_process.configure(state="disabled")

        # 统计变量
        total_files = len(self.ppt_files)  # 总文件数
        success_count = 0  # 成功数量
        fail_count = 0  # 失败数量
        last_output_file = None  # 记录最后一个成功的输出文件

        # 显示开始日志
        self.log_status("\n" + "="*50)
        self.log_status("开始处理...")
        self.log_status("="*50)

        # 遍历每个PPT文件
        # enumerate()同时返回索引和值，从1开始编号
        for idx, ppt_file in enumerate(self.ppt_files, 1):
            try:
                # 获取文件名
                filename = os.path.basename(ppt_file)

                # 更新当前处理文件显示
                self.current_file_label.configure(
                    text=f"📄 正在处理: {filename}",
                    text_color="#1F6AA5"
                )

                # 更新进度条和百分比
                progress = (idx - 1) / total_files  # 计算进度（0-1之间）
                self.progress_bar.set(progress)
                self.progress_label.configure(
                    text=f"文件进度: {idx-1}/{total_files} "
                         f"({progress*100:.0f}%)"
                )

                # 清空页面进度显示
                self.page_progress_label.configure(text="")

                self.log_status(
                    f"\n[{idx}/{total_files}] 正在处理: {filename}"
                )

                # 调用核心处理函数
                output_file = self.replace_background(ppt_file)

                # 处理成功
                output_name = os.path.basename(output_file)
                self.log_status(f"✓ 处理成功: {output_name}")
                success_count += 1
                last_output_file = output_file  # 记录成功的文件路径

            except Exception as e:
                # 处理失败，记录错误
                import traceback
                error_detail = traceback.format_exc()  # 获取详细错误信息
                self.log_status(f"✗ 处理失败: {str(e)}")
                self.log_status(f"详细错误:\n{error_detail}")
                fail_count += 1

        # 所有文件处理完成
        self.progress_bar.set(1.0)  # 进度条设为100%

        # 更新所有进度显示为完成状态
        self.current_file_label.configure(
            text="✅ 全部处理完成！",
            text_color="#2B8A3E"  # 绿色
        )
        self.progress_label.configure(
            text=f"✓ 完成: {total_files}/{total_files} (100%)",
            text_color="#2B8A3E"  # 绿色
        )
        self.page_progress_label.configure(
            text=f"成功: {success_count} 个 | 失败: {fail_count} 个",
            # 如果全部成功显示绿色，否则显示红色
            text_color="#2B8A3E" if fail_count == 0 else "#C92A2A"
        )

        # 显示完成日志
        self.log_status("\n" + "="*50)
        status_msg = f"处理完成！成功: {success_count}, 失败: {fail_count}"
        self.log_status(status_msg)
        self.log_status("="*50 + "\n")

        # 恢复按钮状态
        self.is_processing = False
        self.btn_process.configure(state="normal")

        # 显示完成对话框
        result_msg = (
            f"处理完成！\n\n"
            f"成功: {success_count} 个文件\n"
            f"失败: {fail_count} 个文件"
        )
        messagebox.showinfo("处理完成", result_msg)

        # 如果有成功的文件，自动打开文件所在位置
        if last_output_file and success_count > 0:
            try:
                import subprocess
                # 使用Windows资源管理器打开并选中文件
                # explorer /select, 命令会打开资源管理器并选中指定文件
                subprocess.Popen(
                    f'explorer /select,"{last_output_file}"'
                )
                self.log_status("✓ 已打开文件所在位置")
            except Exception as e:
                # 如果打开失败，记录错误但不影响程序
                self.log_status(f"无法打开文件位置: {str(e)}")

    # ============================================================
    # 第七部分：核心功能 - 替换背景（⭐最重要）
    # ============================================================

    def replace_background(self, ppt_file):
        """
        替换单个PPT的背景
        
        这是程序的核心功能：
        1. 打开PPT文件
        2. 删除原有背景
        3. 添加新背景图片
        4. 保存为新文件
        
        参数:
            ppt_file: PPT文件的路径
            
        返回:
            新文件的路径
        """
        from pptx.util import Inches  # 单位转换（英寸）

        self.log_status("  → 打开PPT文件...")
        
        # 使用python-pptx库打开PPT
        prs = Presentation(ppt_file)
        
        self.log_status(f"  → PPT共有 {len(prs.slides)} 页")

        # 获取幻灯片的宽度和高度（用于设置背景图片大小）
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # 遍历每一页幻灯片
        for idx, slide in enumerate(prs.slides, 1):
            # 更新页面进度显示（实时显示正在处理第几页）
            page_progress = idx / len(prs.slides)
            self.page_progress_label.configure(
                text=f"页面进度: {idx}/{len(prs.slides)} "
                     f"({page_progress*100:.0f}%)"
            )

            self.log_status(f"  → 处理第 {idx}/{len(prs.slides)} 页...")

            # === 步骤1: 删除原有背景 ===
            try:
                # 获取幻灯片的XML元素
                slide_elem = slide.element

                # XML命名空间（用于查找元素）
                namespaces = {
                    'p': 'http://schemas.openxmlformats.org/'
                         'presentationml/2006/main',
                    'a': 'http://schemas.openxmlformats.org/'
                         'drawingml/2006/main'
                }

                # 删除背景填充元素
                bg_removed = False
                for bg in slide_elem.findall('.//p:bg', namespaces):
                    bg.getparent().remove(bg)
                    bg_removed = True

                if bg_removed:
                    self.log_status("    ✓ 已清除背景填充")

                # 隐藏母版背景
                cSld = slide_elem.find('.//p:cSld', namespaces)
                if cSld is not None:
                    cSld.set('showMasterSp', '0')
                    self.log_status("    ✓ 已隐藏母版背景")

                # 删除看起来像背景的大图片和形状（包括渐变蒙版）
                shapes_to_remove = []
                tolerance = Inches(0.15)  # 误差容限
                slide_area = slide_width * slide_height  # 幻灯片面积

                # 检查每个形状
                for shape in slide.shapes:
                    try:
                        # 跳过占位符（如标题框、内容框等）
                        # 占位符是用户输入内容的地方，不应该被删除
                        if (hasattr(shape, 'is_placeholder') and
                                shape.is_placeholder):
                            continue

                        # 检查形状类型：
                        # 1 = 自动形状（如矩形、圆形等，包括渐变蒙版）
                        # 13 = 图片
                        # 6 = 组合形状
                        if shape.shape_type in [1, 13, 6]:
                            # 计算形状的面积和尺寸占比
                            shape_area = shape.width * shape.height
                            area_ratio = shape_area / slide_area  # 面积占比
                            width_ratio = shape.width / slide_width  # 宽度占比
                            height_ratio = shape.height / slide_height  # 高度占比

                            # 检查是否横向撑满（严格匹配）
                            is_width_full = (
                                abs(shape.width - slide_width) < tolerance
                            )

                            # 检查是否纵向撑满（严格匹配）
                            is_height_full = (
                                abs(shape.height - slide_height) < tolerance
                            )

                            # 检查是否接近撑满（多个级别）
                            is_width_near_full = width_ratio > 0.85  # 非常接近
                            is_height_near_full = height_ratio > 0.85
                            is_width_large = width_ratio > 0.70  # 较大（新增）
                            is_height_large = height_ratio > 0.70

                            # 检查是否完全铺满
                            is_full_size = is_width_full and is_height_full

                            # 检查面积占比（多级阈值）
                            is_large_area = area_ratio > 0.6  # 大面积
                            is_medium_area = area_ratio > 0.4  # 中等面积
                            is_small_area = area_ratio > 0.15  # 小面积

                            # 判断是否为背景图片或形状（优化后的规则）
                            # 规则说明（已优化，增加了对渐变蒙版和装饰条的识别）：
                            # 1. 完全铺满（宽高都撑满）-> 删除
                            # 2. 横向撑满且面积>20% -> 删除
                            # 3. 纵向撑满且面积>20% -> 删除
                            # 4. 宽度接近撑满(>85%)且面积>30% -> 删除
                            # 5. 高度接近撑满(>85%)且面积>30% -> 删除
                            # 6. 宽度较大(>70%)且面积>40% -> 删除（捕获底部装饰）
                            # 7. 高度较大(>70%)且面积>40% -> 删除
                            # 8. 宽度很大(>85%)且面积>15% -> 删除（捕获窄装饰条）
                            # 9. 高度很大(>85%)且面积>15% -> 删除
                            # 10. 面积占比>60% -> 删除
                            is_background = (
                                is_full_size or
                                (is_width_full and area_ratio > 0.2) or
                                (is_height_full and area_ratio > 0.2) or
                                (is_width_near_full and area_ratio > 0.3) or
                                (is_height_near_full and area_ratio > 0.3) or
                                (is_width_large and is_medium_area) or
                                (is_height_large and is_medium_area) or
                                (is_width_near_full and is_small_area) or
                                (is_height_near_full and is_small_area) or
                                is_large_area
                            )

                            # 如果是背景图片或形状，标记删除
                            if is_background:
                                shapes_to_remove.append(shape)

                                # 判断形状类型（用于日志显示）
                                shape_type_name = "图片"
                                if shape.shape_type == 1:
                                    shape_type_name = "形状"  # 包括渐变蒙版
                                elif shape.shape_type == 6:
                                    shape_type_name = "组合"

                                reason = ""
                                if is_full_size:
                                    reason = "完全铺满"
                                elif is_width_full:
                                    reason = f"横向撑满({area_ratio*100:.0f}%)"
                                elif is_height_full:
                                    reason = f"纵向撑满({area_ratio*100:.0f}%)"
                                elif is_width_near_full or is_width_large:
                                    reason = (
                                        f"宽度{width_ratio*100:.0f}%"
                                        f"/面积{area_ratio*100:.0f}%"
                                    )
                                elif is_height_near_full or is_height_large:
                                    reason = (
                                        f"高度{height_ratio*100:.0f}%"
                                        f"/面积{area_ratio*100:.0f}%"
                                    )
                                else:
                                    reason = f"面积{area_ratio*100:.0f}%"

                                self.log_status(
                                    f"    → 发现背景{shape_type_name}: {reason} "
                                    f"({shape.width/914400:.1f}x"
                                    f"{shape.height/914400:.1f}英寸)"
                                )
                    except Exception:
                        continue

                # 删除识别出的背景图片和形状
                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

                if shapes_to_remove:
                    self.log_status(
                        f"    ✓ 已删除 {len(shapes_to_remove)} 个背景元素"
                    )

            except Exception as e:
                self.log_status(f"    - 处理背景: {str(e)[:50]}")

            # === 步骤2: 添加新的背景图片 ===
            left = Inches(0)  # 左上角坐标
            top = Inches(0)

            # 添加图片到幻灯片
            pic = slide.shapes.add_picture(
                self.background_image,  # 图片路径
                left, top,  # 位置
                width=slide_width,  # 宽度
                height=slide_height  # 高度
            )

            # 将图片移到最底层
            shapes = slide.shapes._spTree
            shapes.remove(pic._element)  # 先移除
            shapes.insert(2, pic._element)  # 插入到底层位置

        self.log_status("  → 生成输出文件...")
        
        # === 步骤3: 生成输出文件名 ===
        file_path = Path(ppt_file)
        new_name = f"{file_path.stem}_新背景{file_path.suffix}"

        # 根据用户选择的输出目录决定保存位置
        if self.output_dir:
            # 用户指定了输出目录，保存到指定目录
            output_file = Path(self.output_dir) / new_name
        else:
            # 用户没有指定，默认保存到原文件所在目录
            output_file = file_path.parent / new_name

        # 如果文件已存在，添加数字后缀避免覆盖
        counter = 1
        while output_file.exists():
            new_name = f"{file_path.stem}_新背景_{counter}{file_path.suffix}"
            if self.output_dir:
                output_file = Path(self.output_dir) / new_name
            else:
                output_file = file_path.parent / new_name
            counter += 1

        self.log_status("  → 保存文件...")
        
        # === 步骤4: 保存PPT ===
        prs.save(str(output_file))
        
        self.log_status(f"  ✓ 已保存: {output_file.name}")

        # 返回新文件路径
        return str(output_file)


# ============================================================
# 第八部分：程序入口
# ============================================================

def main():
    """
    主函数 - 程序的入口点
    
    创建应用程序实例并运行
    """
    app = PPTBackgroundReplacer()  # 创建应用对象
    app.mainloop()  # 启动主循环（程序会一直运行直到关闭窗口）


# 如果直接运行这个文件（而不是被导入），就执行main函数
if __name__ == "__main__":
    main()


